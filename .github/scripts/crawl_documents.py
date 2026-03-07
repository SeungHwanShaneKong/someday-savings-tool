"""
[ZERO-COST-PIPELINE-2026-03-07] 문서 크롤러 (PDF/PPT)
결혼 관련 공개 보고서/문서를 크롤링하여 텍스트를 추출하고
embed-text Edge Function에 POST하여 벡터 임베딩 생성

실행: python .github/scripts/crawl_documents.py
환경변수: EDGE_FUNCTION_URL, EDGE_FUNCTION_KEY, SUPABASE_SERVICE_ROLE_KEY
"""

import os
import re
import sys
import json
import time
import hashlib
import logging
import tempfile
from pathlib import Path
from typing import Optional
from urllib.robotparser import RobotFileParser

import requests

# Optional PDF/PPT libraries
try:
    import PyPDF2
    HAS_PYPDF2 = True
except ImportError:
    HAS_PYPDF2 = False
    print("Warning: PyPDF2 not installed, PDF parsing disabled")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed, PPT parsing disabled")

# ─── Configuration ───────────────────────────────────────────

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger(__name__)

EDGE_FUNCTION_URL = os.environ.get('EDGE_FUNCTION_URL', '')
EDGE_FUNCTION_KEY = os.environ.get('EDGE_FUNCTION_KEY', '')
SUPABASE_SERVICE_ROLE_KEY = os.environ.get('SUPABASE_SERVICE_ROLE_KEY', '')

# 크롤링 대상 문서 목록
DOCUMENT_SOURCES = [
    {
        'name': '통계청 혼인통계 보고서',
        'url': 'https://kostat.go.kr/cms/stBoardSub/a09/downFile.do',
        'type': 'pdf',
        'category': 'wedding_statistics',
        'region': '전국',
    },
    {
        'name': '한국보건사회연구원 결혼비용 보고서',
        'url': 'https://repository.kihasa.re.kr/bitstream/report',
        'type': 'pdf',
        'category': 'wedding_cost',
        'region': '전국',
    },
]

# 요청 간 딜레이 (초)
REQUEST_DELAY = 3.0

# 최대 텍스트 길이 (임베딩 배치 제한)
MAX_TEXT_LENGTH = 10000

# PII 패턴
PII_PATTERNS = [
    (re.compile(r'\b(01[016789])[-.\s]?(\d{3,4})[-.\s]?(\d{4})\b'), r'\1-****-****'),
    (re.compile(r'\b(02|0[3-6][1-9])[-.\s]?(\d{3,4})[-.\s]?(\d{4})\b'), r'\1-****-****'),
    (re.compile(r'\b(\d{6})[-.\s]?([1-4]\d{6})\b'), r'\1-*******'),
    (re.compile(r'\b([a-zA-Z0-9._%+-]+)@([a-zA-Z0-9.-]+\.[a-zA-Z]{2,})\b'),
     lambda m: f"{m.group(1)[0]}***@{m.group(2)[0]}***.{m.group(2).split('.')[-1]}"),
]

# ─── Utilities ───────────────────────────────────────────────

def mask_pii(text: str) -> str:
    """PII 마스킹 (전화번호, 주민번호, 이메일)"""
    masked = text
    for pattern, replacement in PII_PATTERNS:
        if callable(replacement):
            masked = pattern.sub(replacement, masked)
        else:
            masked = pattern.sub(replacement, masked)
    return masked


def check_robots_txt(url: str) -> bool:
    """robots.txt 준수 확인"""
    try:
        from urllib.parse import urlparse
        parsed = urlparse(url)
        robots_url = f"{parsed.scheme}://{parsed.netloc}/robots.txt"

        rp = RobotFileParser()
        rp.set_url(robots_url)
        rp.read()

        return rp.can_fetch('WeddingSemBot', url)
    except Exception as e:
        logger.warning(f"robots.txt check failed for {url}: {e}")
        return True  # 실패 시 허용


def content_hash(text: str) -> str:
    """SHA-256 content hash for dedup"""
    return hashlib.sha256(text.encode('utf-8')).hexdigest()


def extract_text_from_pdf(filepath: str) -> str:
    """PDF에서 텍스트 추출"""
    if not HAS_PYPDF2:
        logger.warning("PyPDF2 not available, skipping PDF")
        return ''

    try:
        text_parts = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

        return '\n\n'.join(text_parts)
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        return ''


def extract_text_from_pptx(filepath: str) -> str:
    """PPTX에서 텍스트 추출"""
    if not HAS_PPTX:
        logger.warning("python-pptx not available, skipping PPTX")
        return ''

    try:
        prs = Presentation(filepath)
        text_parts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text.strip())

        return '\n\n'.join(text_parts)
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return ''


def download_file(url: str, suffix: str = '.pdf') -> Optional[str]:
    """파일 다운로드 → 임시 파일 경로 반환"""
    try:
        headers = {
            'User-Agent': 'WeddingSemBot/1.0 (wedding cost research)',
            'Accept': 'application/pdf,application/vnd.openxmlformats-officedocument.presentationml.presentation,*/*',
        }

        resp = requests.get(url, headers=headers, timeout=30, stream=True)
        resp.raise_for_status()

        tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
        for chunk in resp.iter_content(chunk_size=8192):
            tmp.write(chunk)
        tmp.close()

        return tmp.name
    except Exception as e:
        logger.error(f"Download failed {url}: {e}")
        return None


def send_to_embed_api(texts: list[str], metadata: list[dict],
                       source_type: str = 'document_crawl',
                       region: str = '') -> bool:
    """embed-text Edge Function에 POST"""
    if not EDGE_FUNCTION_URL or not EDGE_FUNCTION_KEY:
        logger.error("EDGE_FUNCTION_URL or EDGE_FUNCTION_KEY not set")
        return False

    try:
        url = f"{EDGE_FUNCTION_URL}/functions/v1/embed-text"
        headers = {
            'Authorization': f'Bearer {SUPABASE_SERVICE_ROLE_KEY}',
            'Content-Type': 'application/json',
            'apikey': EDGE_FUNCTION_KEY,
        }

        payload = {
            'texts': texts,
            'metadata': metadata,
            'source_type': source_type,
            'region': region or None,
        }

        resp = requests.post(url, json=payload, headers=headers, timeout=60)
        resp.raise_for_status()

        result = resp.json()
        logger.info(f"  Embedded: {result.get('processed', 0)} texts, "
                    f"{result.get('records', 0)} records")
        return True
    except Exception as e:
        logger.error(f"Embed API error: {e}")
        return False


def chunk_text(text: str, max_chars: int = 2048, overlap: int = 200) -> list[str]:
    """텍스트를 임베딩용 청크로 분할"""
    if len(text) <= max_chars:
        return [text] if len(text) > 50 else []

    chunks = []
    start = 0

    while start < len(text):
        end = start + max_chars

        if end < len(text):
            # 문장 경계에서 분할 시도
            last_period = text.rfind('.', start + max_chars // 2, end)
            last_newline = text.rfind('\n', start + max_chars // 2, end)
            break_point = max(last_period, last_newline)

            if break_point > start:
                end = break_point + 1

        chunk = text[start:end].strip()
        if len(chunk) > 50:
            chunks.append(chunk)

        start = end - overlap

    return chunks


# ─── Main Crawl Logic ────────────────────────────────────────

def crawl_source(source: dict) -> dict:
    """단일 소스 크롤링"""
    name = source['name']
    url = source['url']
    doc_type = source.get('type', 'pdf')

    logger.info(f"Processing: {name}")

    # robots.txt 확인
    if not check_robots_txt(url):
        logger.warning(f"  Blocked by robots.txt: {url}")
        return {'name': name, 'status': 'blocked', 'reason': 'robots.txt'}

    # 파일 다운로드
    suffix = '.pdf' if doc_type == 'pdf' else '.pptx'
    filepath = download_file(url, suffix=suffix)

    if not filepath:
        return {'name': name, 'status': 'error', 'reason': 'download_failed'}

    try:
        # 텍스트 추출
        if doc_type == 'pdf':
            text = extract_text_from_pdf(filepath)
        elif doc_type in ('ppt', 'pptx'):
            text = extract_text_from_pptx(filepath)
        else:
            logger.warning(f"  Unknown type: {doc_type}")
            return {'name': name, 'status': 'error', 'reason': f'unknown_type:{doc_type}'}

        if not text or len(text) < 100:
            return {'name': name, 'status': 'skipped', 'reason': 'too_short'}

        # PII 마스킹
        clean_text = mask_pii(text)

        # 길이 제한
        if len(clean_text) > MAX_TEXT_LENGTH:
            clean_text = clean_text[:MAX_TEXT_LENGTH]

        # 청크 분할
        chunks = chunk_text(clean_text)
        if not chunks:
            return {'name': name, 'status': 'skipped', 'reason': 'no_chunks'}

        logger.info(f"  Extracted {len(clean_text)} chars, {len(chunks)} chunks")

        # 메타데이터 구성
        crawled_at = time.strftime('%Y-%m-%dT%H:%M:%SZ', time.gmtime())
        metadata_list = [
            {
                'source': name,
                'url': url,
                'crawled_at': crawled_at,
                'category': source.get('category', 'wedding_cost'),
                'document_type': doc_type,
                'chunk_index': i,
                'total_chunks': len(chunks),
            }
            for i in range(len(chunks))
        ]

        # embed-text API 호출
        success = send_to_embed_api(
            texts=chunks,
            metadata=metadata_list,
            source_type='document_crawl',
            region=source.get('region', ''),
        )

        status = 'success' if success else 'error'
        return {'name': name, 'status': status, 'chunks': len(chunks)}

    finally:
        # 임시 파일 삭제
        try:
            os.unlink(filepath)
        except OSError:
            pass


def main():
    """메인 실행"""
    logger.info("═" * 50)
    logger.info("  Wedding Document Crawler")
    logger.info("  [ZERO-COST-PIPELINE-2026-03-07]")
    logger.info("═" * 50)

    if not EDGE_FUNCTION_URL:
        logger.error("EDGE_FUNCTION_URL environment variable not set")
        sys.exit(1)

    results = {
        'processed': 0,
        'skipped': 0,
        'errors': 0,
        'blocked': 0,
        'details': [],
    }

    for source in DOCUMENT_SOURCES:
        try:
            result = crawl_source(source)
            results['details'].append(result)

            if result['status'] == 'success':
                results['processed'] += 1
            elif result['status'] == 'skipped':
                results['skipped'] += 1
            elif result['status'] == 'blocked':
                results['blocked'] += 1
            else:
                results['errors'] += 1

        except Exception as e:
            logger.error(f"Unexpected error for {source['name']}: {e}")
            results['errors'] += 1
            results['details'].append({
                'name': source['name'],
                'status': 'error',
                'reason': str(e),
            })

        # 요청 간 딜레이
        time.sleep(REQUEST_DELAY)

    # 결과 요약
    logger.info("")
    logger.info("═" * 50)
    logger.info("  Crawl Summary")
    logger.info(f"  Processed: {results['processed']}")
    logger.info(f"  Skipped:   {results['skipped']}")
    logger.info(f"  Blocked:   {results['blocked']}")
    logger.info(f"  Errors:    {results['errors']}")
    logger.info("═" * 50)

    # GitHub Actions output
    print(f"\n::notice::Document crawl complete: "
          f"{results['processed']} processed, "
          f"{results['skipped']} skipped, "
          f"{results['errors']} errors")

    # 에러가 있어도 종료 코드 0 (파이프라인 중단 방지)
    sys.exit(0)


if __name__ == '__main__':
    main()
